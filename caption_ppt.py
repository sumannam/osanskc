from pptx import Presentation
import queue
from contents import *

hymn_file = txt_path + "hymn_list.txt"
hymn_list = open(hymn_file, 'r')
hymn_queue = queue.Queue()

responsive_reading_file = txt_path + "responsive_reading.txt"
responsive_reading_list = open(responsive_reading_file, 'r', encoding='UTF8')

# 찬송가 제목 추출(Queue에 저장)
def createHymnList():
    # hymn_list.txt파일에서 찬송가 제목 끝에 '\n'이 포함되어 있음
    while True:
        hymn = hymn_list.readline()

        if( hymn.find(hymn_file1) == 0):
            hymn = hymn.replace("\n", "")
            hymn_queue.put(hymn)
        if( hymn.find(hymn_file2) == 0):
            hymn = hymn.replace("\n", "")
            hymn_queue.put(hymn)
        if( hymn.find(hymn_file3) == 0):
            hymn = hymn.replace("\n", "")
            hymn_queue.put(hymn)
        if( hymn.find(hymn_file4) == 0):
            hymn = hymn.replace("\n", "")
            hymn_queue.put(hymn)
        if( hymn.find(hymn_file5) == 0):
            hymn = hymn.replace("\n", "")
            hymn_queue.put(hymn)
        if( hymn.find(hymn_file6) == 0):
            hymn = hymn.replace("\n", "")
            hymn_queue.put(hymn)

        if not hymn:
            break

# 찬송가 제목 생성
def makeHymnString(hymn_str):
    length = len(hymn_str)

    if(hymn_str[0:2] == "00"):
        hymn_title = hymn_str[2:3] + "장 " + hymn_str[4:length]
    elif(hymn_str[0] == "0"):
        hymn_title = hymn_str[1:3] + "장 " + hymn_str[4:length]
    else:
        hymn_title = hymn_str[0:3] + "장 " + hymn_str[4:length]
    
    return hymn_title


# 교도문 제목 생성
def makeResponsiveReadingString(num):
    while True:
        responsive_reading_title = responsive_reading_list.readline()

        if( responsive_reading_title.find(num) == 0):
            return responsive_reading_title


createHymnList()

# 파일 열기
prs = Presentation('caption.pptx')


# Slide  1: 축복합니다. 환영합니다.

# Slide  2: 준비찬양 장

# Slide  3:  1. 묵도

# # Slide  4[3]:  2. 찬송1
title_slide_layout = prs.slide_layouts[3]
slide = prs.slides[3]
subtitle = slide.placeholders[1]
subtitle.text = " (찬송) " + makeHymnString(hymn_queue.get())

# # Slide  5[4]:  2. 찬송2
title_slide_layout = prs.slide_layouts[4]
slide = prs.slides[4]
subtitle = slide.placeholders[1]
subtitle.text = " (찬송) " + makeHymnString(hymn_queue.get())

# Slide  6[5]:  3. 성시교독
title_slide_layout = prs.slide_layouts[5]
slide = prs.slides[5]
subtitle = slide.placeholders[1]
subtitle.text = " (교도문) " + makeResponsiveReadingString(responsive_reading)

# Slide  7:  4. 신앙고백

# Slide  8[7]:  5. 찬송
title_slide_layout = prs.slide_layouts[7]
slide = prs.slides[7]
subtitle = slide.placeholders[1]
subtitle.text = " (찬송) " + makeHymnString(hymn_queue.get())

# Slide  9[8]:  6. 기도
title_slide_layout = prs.slide_layouts[8]
slide = prs.slides[8]
subtitle = slide.placeholders[1]
subtitle.text = " (기도) " + prayer

# Slide 10[9]:  7. 성경봉독
title_slide_layout = prs.slide_layouts[9]
slide = prs.slides[9]
subtitle = slide.placeholders[1]
subtitle.text = " (성경봉독) " + short_bible

# Slide 11[10]:  8. 찬송
title_slide_layout = prs.slide_layouts[10]
slide = prs.slides[10]
subtitle = slide.placeholders[1]
subtitle.text = " (찬송) " + makeHymnString(hymn_queue.get())

# Slide 12[11]: 9. 설교
title_slide_layout = prs.slide_layouts[11]
slide = prs.slides[11]
subtitle = slide.placeholders[1]
subtitle.text = " (설교) " + preaching_title

# Slide 13[12]: 10. 헌금
title_slide_layout = prs.slide_layouts[12]
slide = prs.slides[12]
subtitle = slide.placeholders[1]
subtitle.text = " (찬송) " + makeHymnString(hymn_queue.get())

# Slide 14: 11. 헌금기도
# Slide 15: 12. 교회소식

# Slide 16[15]: 13. 찬송
title_slide_layout = prs.slide_layouts[15]
slide = prs.slides[15]
subtitle = slide.placeholders[1]
subtitle.text = " (찬송) " + makeHymnString(hymn_queue.get())

# Slide 17: 14. 축도
# Slide 18: 주님이 함께 하십니다.

prs.save('new-caption.pptx')